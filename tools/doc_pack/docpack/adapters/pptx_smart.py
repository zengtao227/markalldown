from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

from docpack.models import ImageArtifact, MissingDependencyError, PackOptions, PackResult, ProcessingError

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None

try:
    import fitz
except ImportError:  # pragma: no cover - optional dependency
    fitz = None

try:
    from PIL import Image
except ImportError:  # pragma: no cover - optional dependency
    Image = None


def build_pack(options: PackOptions) -> PackResult:
    if Presentation is None:
        raise MissingDependencyError(
            "python-pptx is not installed. Install dependencies from tools/doc_pack/requirements.txt."
        )

    presentation = Presentation(str(options.source_path))
    total_slides = len(presentation.slides)
    selected = options.selected_units or list(range(1, total_slides + 1))
    selected_set = set(selected)

    sections = [f"# Slide Deck Summary\n", f"Source: `{options.source_path.name}`\n"]
    low_text_slides: list[int] = []
    metadata_slides: list[dict[str, object]] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number not in selected_set:
            continue

        text = _extract_slide_text(slide)
        text_length = len(text.replace("\n", " ").strip())
        if text_length < 80:
            low_text_slides.append(slide_number)

        metadata_slides.append(
            {
                "slide": slide_number,
                "text_length": text_length,
                "low_text": slide_number in low_text_slides,
            }
        )
        sections.extend(
            [
                f"## Slide {slide_number}",
                "",
                text or "_No extractable text._",
                "",
            ]
        )

    image_targets = options.image_units or low_text_slides
    images: list[ImageArtifact] = []
    warnings: list[str] = []
    if image_targets:
        try:
            images = _export_slide_images(options, image_targets)
        except (MissingDependencyError, ProcessingError) as exc:
            warnings.append(str(exc))

    if low_text_slides and not images:
        warnings.append(
            "Low-text slides were detected but slide images were not exported. Install LibreOffice to render PPTX -> PDF."
        )

    return PackResult(
        backend="pptx-smart",
        document_type="slides",
        selected_units=selected,
        content_markdown="\n".join(section.strip() for section in sections if section.strip()),
        warnings=warnings,
        filters_applied=[],
        images=images,
        metadata={
            "slide_count": total_slides,
            "slides": metadata_slides,
            "low_text_slides": low_text_slides,
        },
    )


def _extract_slide_text(slide: object) -> str:
    lines: list[str] = []
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        if text and text.strip():
            lines.append(text.strip())
    return "\n\n".join(lines)


def _export_slide_images(options: PackOptions, slide_numbers: list[int]) -> list[ImageArtifact]:
    if fitz is None or Image is None:
        raise MissingDependencyError(
            "PyMuPDF and Pillow are required for slide image export. Install tools/doc_pack/requirements.txt."
        )

    soffice = shutil.which("soffice")
    if not soffice:
        raise ProcessingError("`soffice` was not found. PPTX image export needs LibreOffice.")

    pdf_path = _convert_pptx_to_pdf(soffice, options)
    document = fitz.open(pdf_path)

    artifacts: list[ImageArtifact] = []
    for slide_number in sorted(set(slide_numbers)):
        if slide_number < 1 or slide_number > len(document):
            continue
        page = document.load_page(slide_number - 1)
        artifact = _save_page_image(
            page=page,
            output_dir=options.output_dir / "images",
            prefix="slide",
            unit=slide_number,
            dpi=options.dpi,
            image_format=options.image_format,
        )
        artifacts.append(artifact)

    document.close()
    return artifacts


def _convert_pptx_to_pdf(soffice: str, options: PackOptions) -> Path:
    subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(options.work_dir),
            str(options.source_path),
        ],
        check=True,
        capture_output=True,
        text=True,
        timeout=180,
    )
    pdf_path = options.work_dir / f"{options.source_path.stem}.pdf"
    if not pdf_path.exists():
        raise ProcessingError("LibreOffice conversion did not produce a PDF.")
    return pdf_path


def _save_page_image(
    page: object,
    output_dir: Path,
    prefix: str,
    unit: int,
    dpi: int,
    image_format: str,
) -> ImageArtifact:
    scale = dpi / 72.0
    pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
    suffix = "jpg" if image_format == "jpeg" else "png"
    path = output_dir / f"{prefix}_{unit:03d}_page.{suffix}"

    if image_format == "png":
        pixmap.save(str(path))
    else:
        image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
        image.save(path, format="JPEG", quality=85, optimize=True)

    return ImageArtifact(
        path=path,
        unit=unit,
        kind=prefix,
        width=pixmap.width,
        height=pixmap.height,
    )
